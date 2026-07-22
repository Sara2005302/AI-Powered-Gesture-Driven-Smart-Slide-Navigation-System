import cv2
import os
import time
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from HandTracker import HandDetector
from dottedline import drawrect
import numpy as np
import json
import pyautogui
import traceback
from PIL import Image, ImageDraw, ImageTk

# Import our modules
from modules.config_manager import ConfigManager
from modules.gesture_detector import GestureDetector

class CreativeGestureSystem:
    def __init__(self):
        # Initialize components
        self.config_manager = ConfigManager()
        self.gesture_detector = GestureDetector(self.config_manager)
        
        # System variables
        self.setup_system_variables()
        self.setup_camera()
        self.setup_ui()
    
    def setup_system_variables(self):
        self.width, self.height = 1280, 720
        self.frames_folder = "Images"
        self.slide_num = 0
        
        # Check if Images folder exists
        if os.path.exists(self.frames_folder):
            self.path_imgs = sorted([f for f in os.listdir(self.frames_folder) if f.endswith(('.jpg', '.png'))], 
                                   key=lambda x: int(x.split('.')[0]))
        else:
            print(f"❌ Warning: {self.frames_folder} folder not found!")
            self.path_imgs = []
        
        # Gesture control
        self.gest_done = False
        self.gest_counter = 0
        self.delay = 10
        self.annotations = [[]]
        self.annot_num = 0
        self.annot_start = False
        self.current_color = (0, 0, 255)
        
        # UI control
        self.is_running = False
        self.current_gesture = "Waiting..."
        self.fps = 0
        
        # Camera overlay
        self.hs, self.ws = int(120 * 1.2), int(213 * 1.2)
        self.ge_thresh_y = 400
        self.ge_thresh_x = 750
        
        # System state
        self.cv_window_open = False
        self.safe_exit_flag = False
    
    def setup_camera(self):
        self.cap = cv2.VideoCapture(0)
        if not self.cap.isOpened():
            print("❌ Error: Could not open camera")
            return
        self.cap.set(3, self.width)
        self.cap.set(4, self.height)
        self.detector = HandDetector(detectionCon=0.7, maxHands=1)
    
    def setup_ui(self):
        self.root = tk.Tk()
        self.root.title("✨ AI Gesture Presentation System")
        self.root.geometry("800x600")
        self.root.configure(bg='#0f1b2d')
        
        self.create_title_bar()
        self.create_main_interface()
    
    def create_title_bar(self):
        title_bar = tk.Frame(self.root, bg='#1a2b4a', height=60)
        title_bar.pack(fill="x")
        title_bar.pack_propagate(False)
        
        title_label = tk.Label(title_bar, text="🎭 AI GESTURE PRESENTATION SYSTEM", 
                              font=("Arial", 16, "bold"), fg="#00ff88", bg='#1a2b4a')
        title_label.pack(side="left", padx=20, pady=15)
        
        self.pulse_indicator = tk.Label(title_bar, text="●", font=("Arial", 20), 
                                       fg="red", bg='#1a2b4a')
        self.pulse_indicator.pack(side="right", padx=20)
        self.animate_pulse()
    
    def animate_pulse(self):
        if self.is_running:
            colors = ["green", "lime"]
        else:
            colors = ["red", "darkred"]
            
        current_color = self.pulse_indicator.cget("fg")
        next_color = colors[1] if current_color == colors[0] else colors[0]
        self.pulse_indicator.config(fg=next_color)
        self.root.after(500, self.animate_pulse)
    
    def create_main_interface(self):
        main_frame = tk.Frame(self.root, bg='#0f1b2d')
        main_frame.pack(fill="both", expand=True, padx=20, pady=20)
        
        # Top: Control buttons
        self.create_control_section(main_frame)
        
        # Middle: Status and customization
        middle_frame = tk.Frame(main_frame, bg='#0f1b2d')
        middle_frame.pack(fill="x", pady=10)
        
        self.create_status_section(middle_frame)
        self.create_customization_section(middle_frame)
        
        # Bottom: Gesture guide
        self.create_gesture_guide(main_frame)
    
    def create_control_section(self, parent):
        control_frame = tk.Frame(parent, bg='#152238', relief='raised', bd=2)
        control_frame.pack(fill="x", pady=(0, 10))
        
        # Start/Stop buttons
        btn_frame = tk.Frame(control_frame, bg='#152238')
        btn_frame.pack(pady=15)
        
        self.start_btn = tk.Button(btn_frame, text="🚀 LAUNCH SYSTEM", 
                                  command=self.start_system, bg="#00b894", fg="white",
                                  font=("Arial", 14, "bold"), height=2, width=20)
        self.start_btn.pack(side="left", padx=10)
        
        self.stop_btn = tk.Button(btn_frame, text="🛑 STOP SYSTEM", 
                                 command=self.stop_system, bg="#e17055", fg="white",
                                 font=("Arial", 14, "bold"), height=2, width=20, state="disabled")
        self.stop_btn.pack(side="left", padx=10)
        
        # UPLOAD PPT BUTTON
        self.upload_btn = tk.Button(btn_frame, text="📤 UPLOAD PPT", 
                                   command=self.upload_ppt, bg="#9b59b6", fg="white",
                                   font=("Arial", 14, "bold"), height=2, width=15)
        self.upload_btn.pack(side="left", padx=10)
        
        # Manual controls
        manual_frame = tk.Frame(control_frame, bg='#152238')
        manual_frame.pack(pady=10)
        
        tk.Button(manual_frame, text="⬅️ Previous", command=self.prev_slide,
                 bg="#0984e3", fg="white", font=("Arial", 10), width=10).pack(side="left", padx=5)
        
        tk.Button(manual_frame, text="➡️ Next", command=self.next_slide,
                 bg="#0984e3", fg="white", font=("Arial", 10), width=10).pack(side="left", padx=5)
        
        tk.Button(manual_frame, text="🧹 Clear", command=self.clear_drawings,
                 bg="#fdcb6e", fg="white", font=("Arial", 10), width=10).pack(side="left", padx=5)
        
        # Exit instruction
        exit_label = tk.Label(control_frame, text="💡 Press 'Q' in camera window to close system", 
                             font=("Arial", 10), fg="#ffd700", bg='#152238')
        exit_label.pack(pady=5)
    
    def create_status_section(self, parent):
        status_frame = tk.Frame(parent, bg='#152238', relief='raised', bd=2)
        status_frame.pack(side="left", fill="both", expand=True, padx=(0, 10))
        
        # Current gesture
        gesture_frame = tk.Frame(status_frame, bg='#152238')
        gesture_frame.pack(pady=10)
        
        self.gesture_display = tk.Label(gesture_frame, text="👋", font=("Arial", 24), 
                                       bg='#152238', fg="white")
        self.gesture_display.pack()
        
        self.gesture_text = tk.Label(gesture_frame, text="Ready to Detect", 
                                   font=("Arial", 12, "bold"), bg='#152238', fg="#00ff88")
        self.gesture_text.pack()
        
        # Slide info
        slide_frame = tk.Frame(status_frame, bg='#152238')
        slide_frame.pack(pady=10)
        
        self.slide_label = tk.Label(slide_frame, text="No Slides Loaded", 
                                  font=("Arial", 10), bg='#152238', fg="white")
        self.slide_label.pack()
        
        self.fps_label = tk.Label(slide_frame, text="FPS: 0", 
                                 font=("Arial", 10), bg='#152238', fg="white")
        self.fps_label.pack()
    
    def create_customization_section(self, parent):
        custom_frame = tk.Frame(parent, bg='#152238', relief='raised', bd=2)
        custom_frame.pack(side="right", fill="both", expand=True)
        
        tk.Label(custom_frame, text="⚙️ CUSTOMIZE GESTURES", 
                font=("Arial", 12, "bold"), fg="#ff9800", bg='#152238').pack(pady=10)
        
        # Action selection
        action_frame = tk.Frame(custom_frame, bg='#152238')
        action_frame.pack(fill="x", padx=10, pady=5)
        
        tk.Label(action_frame, text="Action:", font=("Arial", 10), 
                fg="white", bg='#152238').pack(side="left")
        
        self.action_var = tk.StringVar(value="next_slide")
        actions = ["next_slide", "previous_slide", "pointer", "draw", "erase", "clear_all"]
        action_dropdown = ttk.Combobox(action_frame, textvariable=self.action_var,
                                      values=actions, state="readonly", width=12)
        action_dropdown.pack(side="left", padx=10)
        action_dropdown.bind('<<ComboboxSelected>>', self.load_current_gesture)
        
        # Finger selection
        finger_frame = tk.Frame(custom_frame, bg='#152238')
        finger_frame.pack(pady=5)
        
        tk.Label(finger_frame, text="Fingers:", font=("Arial", 10), 
                fg="white", bg='#152238').pack()
        
        self.finger_vars = []
        finger_names = ["Thumb", "Index", "Middle", "Ring", "Pinky"]
        for name in finger_names:
            var = tk.BooleanVar()
            self.finger_vars.append(var)
            cb = tk.Checkbutton(finger_frame, text=name, variable=var,
                               bg='#152238', fg="white", selectcolor="#1a2b4a")
            cb.pack(side="left", padx=2)
        
        # Control buttons
        btn_frame = tk.Frame(custom_frame, bg='#152238')
        btn_frame.pack(pady=10)
        
        tk.Button(btn_frame, text="💾 Save", command=self.save_custom_gesture,
                 bg="#4caf50", fg="white", font=("Arial", 9)).pack(side="left", padx=5)
        
        tk.Button(btn_frame, text="🔄 Reset", command=self.reset_gestures,
                 bg="#f44336", fg="white", font=("Arial", 9)).pack(side="left", padx=5)
        
        self.load_current_gesture()
    
    def create_gesture_guide(self, parent):
        guide_frame = tk.Frame(parent, bg='#152238', relief='raised', bd=2)
        guide_frame.pack(fill="x", pady=(10, 0))
        
        tk.Label(guide_frame, text="🎯 GESTURE GUIDE", 
                font=("Arial", 12, "bold"), fg="#ffd700", bg='#152238').pack(pady=10)
        
        self.update_gesture_guide(guide_frame)
    
    def update_gesture_guide(self, parent):
        """Create gesture guide with current mappings"""
        # Clear existing guide
        for widget in parent.winfo_children():
            if widget != parent.winfo_children()[0]:  # Keep title
                widget.destroy()
        
        guide_container = tk.Frame(parent, bg='#152238')
        guide_container.pack(fill="x", padx=10, pady=5)
        
        # Create 2-column layout
        left_column = tk.Frame(guide_container, bg='#152238')
        left_column.pack(side="left", fill="x", expand=True)
        
        right_column = tk.Frame(guide_container, bg='#152238')
        right_column.pack(side="right", fill="x", expand=True)
        
        gesture_descriptions = {
            'next_slide': 'Next Slide ',
            'previous_slide': 'Previous Slide',
            'pointer': 'Pointer',
            'draw': 'Draw ',
            'erase': 'Erase ',
            'clear_all': 'Clear All ',
            'swipe_right': 'Next Slide ',
            'swipe_left': 'Prev Slide '
        }
        
        gestures = list(gesture_descriptions.items())
        mid_point = len(gestures) // 2
        
        for i, (action, description) in enumerate(gestures):
            column = left_column if i < mid_point else right_column
            
            if action in ['swipe_right', 'swipe_left']:
                # Swipe gestures with fixed icons
                icon_map = {'swipe_right': '➡️', 'swipe_left': '⬅️'}
                icon = icon_map.get(action, '❓')
            else:
                # Customizable gestures
                icon = self.config_manager.get_gesture_icon(action)
            
            gest_frame = tk.Frame(column, bg='#1a2b4a')
            gest_frame.pack(fill="x", pady=2)
            
            tk.Label(gest_frame, text=icon, font=("Arial", 12), 
                   bg='#1a2b4a', fg="#00ff88", width=3).pack(side="left")
            
            tk.Label(gest_frame, text=description, font=("Arial", 9), 
                   bg='#1a2b4a', fg="white", anchor="w").pack(side="left", fill="x", expand=True)
    
    def load_current_gesture(self, event=None):
        """Load current gesture into checkboxes"""
        action = self.action_var.get()
        if action in self.config_manager.gesture_mappings:
            gesture_type = self.config_manager.gesture_mappings[action]
            
            # Reset all checkboxes
            for var in self.finger_vars:
                var.set(False)
            
            # If it's a finger-based gesture, set the checkboxes
            if isinstance(gesture_type, list):
                for i, var in enumerate(self.finger_vars):
                    if i < len(gesture_type):
                        var.set(gesture_type[i] == 1)
    
    def save_custom_gesture(self):
        """Save the custom gesture mapping"""
        action = self.action_var.get()
        fingers = [1 if var.get() else 0 for var in self.finger_vars]
        
        # Update gesture mappings
        self.config_manager.gesture_mappings[action] = fingers
        
        if self.config_manager.save_mappings_to_file():
            print(f"✅ Saved {action} = {fingers}")
            # Update gesture guide
            guide_frame = self.root.winfo_children()[1].winfo_children()[-1]
            self.update_gesture_guide(guide_frame)
        else:
            print(f"❌ Failed to save {action}")
    
    def reset_gestures(self):
        """Reset all gestures to defaults"""
        self.config_manager.gesture_mappings = {
            'next_slide': [0, 0, 0, 0, 1],
            'previous_slide': [1, 0, 0, 0, 0],
            'pointer': [0, 1, 1, 0, 0],
            'draw': [0, 1, 0, 0, 0],
            'erase': [0, 1, 1, 1, 0],
            'clear_all': [1, 1, 1, 1, 1]
        }
        
        if self.config_manager.save_mappings_to_file():
            # Update gesture guide
            guide_frame = self.root.winfo_children()[1].winfo_children()[-1]
            self.update_gesture_guide(guide_frame)
            self.load_current_gesture()
            print("🔄 All gestures reset to defaults")
        else:
            print("❌ Failed to reset gestures")

    # ========== PPT UPLOAD WITH SPIRE.PRESENTATION ==========
    
    def upload_ppt(self):
        """Handle PPT upload and conversion using Spire.Presentation for perfect formatting"""
        # Ask user to select PPT file
        file_path = filedialog.askopenfilename(
            title="Select PowerPoint File",
            filetypes=[("PowerPoint Files", "*.ppt *.pptx"), ("All Files", "*.*")]
        )
        
        if not file_path:
            return
        
        ppt_filename = os.path.basename(file_path)
        print(f"📤 Uploading: {ppt_filename}")
        
        # Show uploading status
        self.gesture_text.config(text="Uploading PPT...", fg="#ff9800")
        self.gesture_display.config(text="📤")
        self.root.update()
        
        try:
            # ✅ Create Images folder if it doesn't exist, but DON'T delete existing images
            if not os.path.exists(self.frames_folder):
                os.makedirs(self.frames_folder)
            
            # ✅ Find the next available slide number (preserve existing images)
            existing_images = [f for f in os.listdir(self.frames_folder) if f.endswith(('.jpg', '.png'))]
            if existing_images:
                existing_numbers = []
                for img in existing_images:
                    try:
                        num = int(img.split('.')[0])
                        existing_numbers.append(num)
                    except ValueError:
                        continue
                start_slide_num = max(existing_numbers) + 1 if existing_numbers else 1
            else:
                start_slide_num = 1
            
            # ✅ Use Spire.Presentation for perfect formatting preservation
            from spire.presentation import Presentation
            
            # Create presentation object
            presentation = Presentation()
            
            # Load the PowerPoint file
            presentation.LoadFromFile(file_path)
            
            slide_count = len(presentation.Slides)
            
            print(f"🔄 Converting {slide_count} slides using Spire.Presentation...")
            self.gesture_text.config(text=f"Converting {slide_count} slides...")
            self.root.update()
            
            # Export each slide as PNG image
            for i, slide in enumerate(presentation.Slides):
                # ✅ Use unique slide numbers that continue from existing slides
                slide_number = start_slide_num + i
                
                # Create image file path
                image_path = os.path.join(self.frames_folder, f"{slide_number}.png")
                
                # Save slide as image (preserves all formatting, images, charts)
                image_stream = slide.SaveAsImage()
                image_stream.Save(image_path)
                
                # Update progress
                progress = f"Converting... {i+1}/{slide_count}"
                self.gesture_text.config(text=progress)
                self.root.update()
            
            # Close presentation to free resources
            presentation.Dispose()
            
            # ✅ Reload ALL slides (existing + new)
            self.path_imgs = sorted([f for f in os.listdir(self.frames_folder) if f.endswith(('.jpg', '.png'))], 
                                   key=lambda x: int(x.split('.')[0]))
            
            # Start from the first new slide
            self.slide_num = start_slide_num - 1
            
            # Update UI
            self.update_slide_display()
            self.gesture_text.config(text="PPT Converted!", fg="#00ff88")
            self.gesture_display.config(text="✅")
            
            total_slides = len(self.path_imgs)
            print(f"✅ Conversion complete! Added {slide_count} slides. Total slides: {total_slides}")
            
            # Auto-start gesture system if not running
            if not self.is_running:
                start_now = messagebox.askyesno(
                    "PPT Converted", 
                    f"Successfully converted {slide_count} slides from {ppt_filename}!\nTotal slides: {total_slides}\n\nStart gesture control now?"
                )
                if start_now:
                    self.start_system()
        
        except ImportError:
            # Spire.Presentation not available, use fallback method
            print("❌ Spire.Presentation not installed, using fallback method")
            self.upload_ppt_fallback(file_path, ppt_filename)
        except Exception as e:
            print(f"❌ Spire.Presentation conversion failed: {e}")
            # Try fallback method
            self.upload_ppt_fallback(file_path, ppt_filename)
    
    def upload_ppt_fallback(self, file_path, filename):
        """Fallback method if Spire.Presentation is not available"""
        try:
            # Try using python-pptx as fallback
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slide_count = len(prs.slides)
            
            print(f"🔄 Using fallback method for {slide_count} slides...")
            self.gesture_text.config(text=f"Converting {slide_count} slides (fallback)...")
            self.root.update()
            
            # Find next available slide number
            existing_images = [f for f in os.listdir(self.frames_folder) if f.endswith(('.jpg', '.png'))]
            if existing_images:
                existing_numbers = []
                for img in existing_images:
                    try:
                        num = int(img.split('.')[0])
                        existing_numbers.append(num)
                    except ValueError:
                        continue
                start_slide_num = max(existing_numbers) + 1 if existing_numbers else 1
            else:
                start_slide_num = 1
            
            for i, slide in enumerate(prs.slides):
                slide_number = start_slide_num + i
                
                # Create basic slide image
                img = Image.new('RGB', (1920, 1080), color='white')
                draw = ImageDraw.Draw(img)
                
                # Add professional background
                self.draw_slide_background(draw, i)
                
                # Add slide content
                draw.text((100, 100), f"Slide {slide_number}", fill='#1565c0', size=80)
                draw.text((100, 200), f"From: {filename}", fill='#d32f2f', size=30)
                
                # Add content from slide
                y_pos = 300
                content_added = False
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text[:100] + "..." if len(shape.text) > 100 else shape.text
                        draw.text((100, y_pos), f"• {text}", fill='#333333', size=24)
                        y_pos += 40
                        content_added = True
                
                if not content_added:
                    draw.text((100, 350), "No text content detected", fill='#666666', size=24)
                    draw.text((100, 400), "Using basic slide layout", fill='#666666', size=20)
                
                # Save image
                img_path = os.path.join(self.frames_folder, f"{slide_number}.jpg")
                img.save(img_path, 'JPEG', quality=95)
                
                # Update progress
                progress = f"Converting... {i+1}/{slide_count}"
                self.gesture_text.config(text=progress)
                self.root.update()
            
            # Reload all slides
            self.path_imgs = sorted([f for f in os.listdir(self.frames_folder) if f.endswith(('.jpg', '.png'))], 
                                   key=lambda x: int(x.split('.')[0]))
            self.slide_num = start_slide_num - 1
            self.update_slide_display()
            
            self.gesture_text.config(text="PPT Converted (Fallback)", fg="#ff9800")
            self.gesture_display.config(text="📊")
            
            print(f"✅ Fallback conversion complete for: {filename}")
            
        except Exception as e:
            print(f"❌ Fallback conversion also failed: {e}")
            messagebox.showerror("Conversion Error", 
                               f"Could not convert PPT file:\n{str(e)}\n\nPlease install Spire.Presentation for better results.")
    
    def draw_slide_background(self, draw, slide_index):
        """Draw professional slide background for fallback method"""
        colors = [
            ('#e3f2fd', '#bbdefb'),  # Light blue gradient
            ('#f3e5f5', '#e1bee7'),  # Light purple gradient  
            ('#e8f5e8', '#c8e6c9'),  # Light green gradient
            ('#fff3e0', '#ffe0b2'),  # Light orange gradient
            ('#fce4ec', '#f8bbd9')   # Light pink gradient
        ]
        
        bg_color, header_color = colors[slide_index % len(colors)]
        
        # Draw gradient background (simplified)
        draw.rectangle([0, 0, 1920, 1080], fill=bg_color)
        
        # Add header bar
        draw.rectangle([0, 0, 1920, 120], fill=header_color)

    # ========== EXISTING METHODS ==========
    
    def start_system(self):
        if not self.path_imgs:
            print("❌ No slides found in Images folder!")
            return
            
        self.is_running = True
        self.start_btn.config(state="disabled", bg="#7f8c8d")
        self.stop_btn.config(state="normal", bg="#e74c3c")
        self.gesture_text.config(text="System Active", fg="#00ff88")
        self.update_slide_display()
        
        # Start gesture processing thread
        self.gesture_thread = threading.Thread(target=self.process_gestures)
        self.gesture_thread.daemon = True
        self.gesture_thread.start()
        
        print("🚀 Gesture system started!")
    
    def stop_system(self):
        self.is_running = False
        self.safe_exit_flag = True
        self.start_btn.config(state="normal", bg="#00b894")
        self.stop_btn.config(state="disabled", bg="#7f8c8d")
        self.gesture_text.config(text="System Stopped", fg="#ff6b6b")
        self.gesture_display.config(text="👋")
        
        # Close OpenCV window if open
        if self.cv_window_open:
            cv2.destroyAllWindows()
            self.cv_window_open = False
        
        print("🛑 Gesture system stopped!")
    
    def prev_slide(self):
        if self.slide_num > 0:
            self.slide_num -= 1
            self.annotations = [[]]
            self.update_slide_display()
    
    def next_slide(self):
        if self.slide_num < len(self.path_imgs) - 1:
            self.slide_num += 1
            self.annotations = [[]]
            self.update_slide_display()
    
    def clear_drawings(self):
        self.annotations = [[]]
        self.annot_num = 0
        print("🧹 Drawings cleared!")
    
    def update_slide_display(self):
        if self.path_imgs:
            self.slide_label.config(text="Slide {}/{}".format(self.slide_num + 1, len(self.path_imgs)))
        else:
            self.slide_label.config(text="No Slides Loaded")
    
    def process_gestures(self):
        if not self.path_imgs:
            print("❌ No slides to display!")
            return
            
        prev_time = 0
        self.cv_window_open = True
        self.safe_exit_flag = False
        
        while self.is_running and not self.safe_exit_flag:
            try:
                current_time = time.time()
                # Calculate FPS
                self.fps = 1 / (current_time - prev_time) if prev_time != 0 else 0
                prev_time = current_time
                
                # Read camera frame
                success, frame = self.cap.read()
                if not success:
                    continue
                
                frame = cv2.flip(frame, 1)
                
                # Load current slide
                pathFullImage = os.path.join(self.frames_folder, self.path_imgs[self.slide_num])
                slide_current = cv2.imread(pathFullImage)
                if slide_current is None:
                    continue
                slide_current = cv2.resize(slide_current, (self.width, self.height))
                
                # Find hands
                hands, frame = self.detector.findHands(frame)
                drawrect(frame, (self.width, 0), (self.ge_thresh_x, self.ge_thresh_y), (0, 255, 0), 5, 'dotted')
                
                current_gesture = "None"
                gesture_icon = "👋"
                
                if hands and len(hands) > 0 and not self.gest_done:
                    hand = hands[0]
                    
                    if not all(key in hand for key in ['center', 'lmList']):
                        continue
                    
                    cx, cy = hand["center"]
                    lm_list = hand["lmList"]
                    
                    if not lm_list or len(lm_list) < 21:
                        continue
                    
                    try:
                        fingers = self.detector.fingersUp(hand)
                        if len(fingers) != 5:
                            fingers = [0, 0, 0, 0, 0]
                    except Exception as e:
                        fingers = [0, 0, 0, 0, 0]
                    
                    if len(lm_list) < 9:
                        continue
                    
                    try:
                        x_val = int(np.interp(lm_list[8][0], [self.width//2, self.width], [0, self.width]))
                        y_val = int(np.interp(lm_list[8][1], [150, self.height - 150], [0, self.height]))
                        index_fing = (x_val, y_val)
                    except Exception as e:
                        index_fing = (0, 0)
                    
                    if cy < self.ge_thresh_y and cx > self.ge_thresh_x:
                        self.annot_start = False
                        
                        # BASIC GESTURE RECOGNITION
                        if fingers == self.config_manager.gesture_mappings['previous_slide']:
                            current_gesture = "Previous Slide"
                            gesture_icon = "👍"
                            if self.slide_num > 0:
                                self.gest_done = True
                                self.slide_num -= 1
                                self.annotations = [[]]
                                self.annot_num = 0
                                self.update_slide_display()
                        
                        elif fingers == self.config_manager.gesture_mappings['next_slide']:
                            current_gesture = "Next Slide"
                            gesture_icon = "🖖"
                            if self.slide_num < len(self.path_imgs) - 1:
                                self.gest_done = True
                                self.slide_num += 1
                                self.annotations = [[]]
                                self.annot_num = 0
                                self.update_slide_display()
                        
                        elif fingers == self.config_manager.gesture_mappings['clear_all']:
                            current_gesture = "Clear All"
                            gesture_icon = "🖐️"
                            if self.annotations:
                                self.annotations.clear()
                                self.annot_num = 0
                                self.gest_done = True
                                self.annotations = [[]]
                    
                    # SWIPE GESTURES
                    try:
                        swipe_gesture, swipe_icon = self.gesture_detector.process_swipe_gestures(hands, cx, cy)
                        if swipe_gesture:
                            current_gesture, gesture_icon = swipe_gesture, swipe_icon
                            
                            if swipe_gesture == "Next Slide (Swipe)":
                                if self.slide_num < len(self.path_imgs) - 1:
                                    self.slide_num += 1
                                    self.annotations = [[]]
                                    self.update_slide_display()
                                    self.gest_done = True
                            elif swipe_gesture == "Previous Slide (Swipe)":
                                if self.slide_num > 0:
                                    self.slide_num -= 1
                                    self.annotations = [[]]
                                    self.update_slide_display()
                                    self.gest_done = True
                    except Exception as e:
                        print(f"⚠️ Swipe detection error: {e}")
                    
                    # Pointer gesture
                    if fingers == self.config_manager.gesture_mappings['pointer']:
                        current_gesture = "Pointer"
                        gesture_icon = "✌️"
                        try:
                            cv2.circle(slide_current, index_fing, 4, (0, 0, 255), cv2.FILLED)
                        except:
                            pass
                        self.annot_start = False
                    
                    # Draw gesture
                    if fingers == self.config_manager.gesture_mappings['draw']:
                        current_gesture = "Drawing"
                        gesture_icon = "👆"
                        if not self.annot_start:
                            self.annot_start = True
                            self.annot_num += 1
                            self.annotations.append([])
                        if self.annot_start:
                            try:
                                self.annotations[self.annot_num].append(index_fing)
                                cv2.circle(slide_current, index_fing, 4, self.current_color, cv2.FILLED)
                            except:
                                pass
                    else:
                        self.annot_start = False
                    
                    # Erase gesture
                    if fingers == self.config_manager.gesture_mappings['erase']:
                        current_gesture = "Erase"
                        gesture_icon = "🤟"
                        if self.annotations and len(self.annotations) > 0:
                            try:
                                self.annotations.pop(-1)
                                self.annot_num = max(0, self.annot_num - 1)
                                self.gest_done = True
                            except:
                                pass
                
                else:
                    self.annot_start = False
                
                # Gesture cooldown
                if self.gest_done:
                    self.gest_counter += 1
                    if self.gest_counter > self.delay:
                        self.gest_counter = 0
                        self.gest_done = False
                
                # Draw annotations
                try:
                    for i, annotation in enumerate(self.annotations):
                        for j in range(len(annotation)):
                            if j != 0:
                                cv2.line(slide_current, annotation[j - 1], annotation[j], self.current_color, 6)
                except Exception as e:
                    print(f"⚠️ Annotation drawing error: {e}")
                
                # Add camera overlay
                try:
                    img_small = cv2.resize(frame, (self.ws, self.hs))
                    h, w, _ = slide_current.shape
                    slide_current[h-self.hs:h, w-self.ws:w] = img_small
                except Exception as e:
                    print(f"⚠️ Camera overlay error: {e}")
                
                # Update UI in main thread
                self.root.after(0, self.update_ui_status, current_gesture, gesture_icon, self.fps)
                
                # Display
                try:
                    cv2.imshow("✨ AI Gesture Presentation - PRESS 'Q' TO CLOSE", slide_current)
                except Exception as e:
                    print(f"⚠️ Display error: {e}")
                    break
                
                # Safe exit detection
                key = cv2.waitKey(1) & 0xFF
                if key == ord('q'):
                    print("🛑 Safe exit initiated by 'Q' key")
                    self.safe_exit_flag = True
                    self.root.after(0, self.stop_system)
                    break
                
                time.sleep(0.01)
                
            except Exception as e:
                print(f"⚠️ Main gesture processing error: {e}")
                continue
        
        # Safe cleanup
        try:
            cv2.destroyAllWindows()
        except:
            pass
        self.cv_window_open = False
        print("✅ Gesture processing stopped safely")
    
    def update_ui_status(self, gesture, icon, fps):
        if self.is_running:
            self.gesture_display.config(text=icon)
            self.gesture_text.config(text=gesture)
            self.fps_label.config(text="FPS: {:.1f}".format(fps))
    
    def run(self):
        try:
            self.root.mainloop()
        finally:
            self.is_running = False
            self.safe_exit_flag = True
            if hasattr(self, 'cap') and self.cap:
                self.cap.release()
            try:
                cv2.destroyAllWindows()
            except:
                pass

# This allows the file to be imported as a module
if __name__ == "__main__":
    print("🚀 Starting AI Gesture Presentation System...")
    system = CreativeGestureSystem()
    system.run()